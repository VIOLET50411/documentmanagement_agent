"""PPTX parser with python-pptx support and XML fallback."""

from __future__ import annotations

from pathlib import Path
from xml.etree import ElementTree
from zipfile import ZipFile

PPTX_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
PPTX_SLIDE_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
RELS_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
DEFAULT_SLIDE_TITLE = "\u7b2c{index}\u9875\u5e7b\u706f\u7247"


class PptxParser:
    """Parse PPTX files into heading, paragraph, and table elements."""

    def parse(self, file_path: str) -> list[dict]:
        path = Path(file_path)

        try:
            elements = self._parse_with_python_pptx(path)
            if elements:
                return elements
        except (OSError, RuntimeError, ValueError, TypeError, KeyError):
            pass

        return self._parse_with_xml(path)

    def _parse_with_python_pptx(self, path: Path) -> list[dict]:
        try:
            from pptx import Presentation
            from pptx.util import Inches  # noqa: F401 - confirms pptx is importable
        except ImportError:
            return []

        prs = Presentation(str(path))
        elements: list[dict] = []

        for slide_index, slide in enumerate(prs.slides, start=1):
            slide_title = self._extract_slide_title(slide, slide_index)
            text_count = 0
            table_count = 0
            slide_items: list[dict] = []

            for shape_index, shape in enumerate(slide.shapes, start=1):
                if shape.has_table:
                    rows = self._extract_table_rows_pptx(shape.table)
                    if rows:
                        table_count += 1
                        slide_items.append(
                            self._build_table_summary_element(
                                rows=rows,
                                section_title=slide_title,
                                slide_index=slide_index,
                                block_index=shape_index,
                                parser="python-pptx",
                            )
                        )
                        slide_items.append(
                            {
                                "type": "table",
                                "text": self._render_markdown_table(rows),
                                "metadata": {
                                    "section_title": slide_title,
                                    "page_number": slide_index,
                                    "block_index": shape_index,
                                    "parser": "python-pptx",
                                },
                            }
                        )
                    continue

                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = (para.text or "").strip()
                        if not text:
                            continue
                        text_count += 1
                        is_heading = para.level == 0 and hasattr(para, "font") and (
                            (para.font and para.font.size and para.font.size.pt >= 20)
                            if para.font
                            else False
                        )
                        slide_items.append(
                            {
                                "type": "heading" if is_heading else "paragraph",
                                "text": text,
                                "metadata": {
                                    "section_title": slide_title,
                                    "page_number": slide_index,
                                    "block_index": shape_index,
                                    "parser": "python-pptx",
                                },
                            }
                        )

            if slide_items:
                elements.append(
                    self._build_slide_overview_element(
                        slide_title=slide_title,
                        slide_index=slide_index,
                        text_count=text_count,
                        table_count=table_count,
                    )
                )
                elements.extend(slide_items)

        return elements

    def _extract_slide_title(self, slide, slide_index: int) -> str:
        if slide.shapes.title and slide.shapes.title.text:
            return slide.shapes.title.text.strip()
        for shape in slide.shapes:
            if hasattr(shape, "name") and "title" in (shape.name or "").lower():
                if shape.has_text_frame and shape.text_frame.text.strip():
                    return shape.text_frame.text.strip()
        return DEFAULT_SLIDE_TITLE.format(index=slide_index)

    def _extract_table_rows_pptx(self, table) -> list[list[str]]:
        rows: list[list[str]] = []
        for row in table.rows:
            cells = [(cell.text or "").strip() for cell in row.cells]
            if any(cells):
                rows.append(cells)
        return rows

    def _parse_with_xml(self, path: Path) -> list[dict]:
        elements: list[dict] = []

        with ZipFile(path) as archive:
            slide_names = sorted(
                [n for n in archive.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")],
                key=lambda n: int("".join(filter(str.isdigit, n.split("/")[-1])) or 0),
            )

            for slide_index, slide_name in enumerate(slide_names, start=1):
                slide_title = DEFAULT_SLIDE_TITLE.format(index=slide_index)
                root = ElementTree.fromstring(archive.read(slide_name))

                texts: list[str] = []
                for paragraph in root.iter(f"{{{PPTX_NS}}}p"):
                    runs = paragraph.findall(f".//{{{PPTX_NS}}}r/{{{PPTX_NS}}}t")
                    text = "".join((r.text or "") for r in runs).strip()
                    if text:
                        texts.append(text)

                slide_items: list[dict] = []
                if texts:
                    slide_title = texts[0]
                    slide_items.append(
                        {
                            "type": "heading",
                            "text": texts[0],
                            "metadata": {
                                "section_title": slide_title,
                                "page_number": slide_index,
                                "block_index": 1,
                                "parser": "xml",
                            },
                        }
                    )
                    for bi, text in enumerate(texts[1:], start=2):
                        slide_items.append(
                            {
                                "type": "paragraph",
                                "text": text,
                                "metadata": {
                                    "section_title": slide_title,
                                    "page_number": slide_index,
                                    "block_index": bi,
                                    "parser": "xml",
                                },
                            }
                        )

                table_index = 0
                for table_node in root.iter(f"{{{PPTX_NS}}}tbl"):
                    rows = self._extract_table_rows_xml(table_node)
                    if rows:
                        table_index += 1
                        slide_items.append(
                            self._build_table_summary_element(
                                rows=rows,
                                section_title=slide_title,
                                slide_index=slide_index,
                                block_index=table_index,
                                parser="xml",
                            )
                        )
                        slide_items.append(
                            {
                                "type": "table",
                                "text": self._render_markdown_table(rows),
                                "metadata": {
                                    "section_title": slide_title,
                                    "page_number": slide_index,
                                    "table_index": table_index,
                                    "parser": "xml",
                                },
                            }
                        )

                if slide_items:
                    elements.append(
                        self._build_slide_overview_element(
                            slide_title=slide_title,
                            slide_index=slide_index,
                            text_count=len(texts),
                            table_count=table_index,
                        )
                    )
                    elements.extend(slide_items)

        return elements

    def _extract_table_rows_xml(self, table_node) -> list[list[str]]:
        rows: list[list[str]] = []
        for row in table_node.iter(f"{{{PPTX_NS}}}tr"):
            cells: list[str] = []
            for cell in row.iter(f"{{{PPTX_NS}}}tc"):
                runs = cell.iter(f"{{{PPTX_NS}}}t")
                cells.append("".join((r.text or "") for r in runs).strip())
            if any(cells):
                rows.append(cells)
        return rows

    def _render_markdown_table(self, rows: list[list[str]]) -> str:
        if not rows:
            return ""
        header = rows[0]
        body = rows[1:]
        rendered = [
            "| " + " | ".join(header) + " |",
            "| " + " | ".join(["---"] * len(header)) + " |",
        ]
        for row in body[:200]:
            padded = row + [""] * max(0, len(header) - len(row))
            rendered.append("| " + " | ".join(padded[: len(header)]) + " |")
        return "\n".join(rendered)

    def _build_slide_overview_element(
        self,
        *,
        slide_title: str,
        slide_index: int,
        text_count: int,
        table_count: int,
    ) -> dict:
        return {
            "type": "paragraph",
            "text": (
                f"\u7b2c{slide_index}\u9875\u5e7b\u706f\u7247\u300a{slide_title}\u300b\u6982\u89c8\uff1a"
                f"\u5305\u542b{text_count}\u6bb5\u6587\u672c\u3001{table_count}\u4e2a\u8868\u683c\u3002"
            ),
            "metadata": {
                "section_title": slide_title,
                "page_number": slide_index,
                "block_index": 0,
                "parser": "pptx_overview",
            },
        }

    def _build_table_summary_element(
        self,
        *,
        rows: list[list[str]],
        section_title: str,
        slide_index: int,
        block_index: int,
        parser: str,
    ) -> dict:
        headers = [cell for cell in rows[0] if cell][:5]
        row_count = max(len(rows) - 1, 0)
        header_text = "\u3001".join(headers) if headers else "\u672a\u8bc6\u522b\u5217\u540d"
        return {
            "type": "paragraph",
            "text": (
                f"\u8868\u683c\u6458\u8981\uff1a\u5171{row_count}\u884c\uff0c"
                f"\u5b57\u6bb5\u5305\u62ec{header_text}\u3002"
            ),
            "metadata": {
                "section_title": section_title,
                "page_number": slide_index,
                "block_index": block_index,
                "parser": f"{parser}_table_summary",
            },
        }
